"""
File Format Converter Module
Handles conversion between PDF and various file formats.
Supports: docx, doc, txt, html, md, jpg, png, ppt/pptx, xls/xlsx → PDF
          PDF → docx, txt, html, md
"""
import os
import re
import logging
import shutil
import subprocess

import docx
import html2text
import markdown
from reportlab.lib.pagesizes import letter
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.platypus import SimpleDocTemplate, Paragraph

logging.basicConfig(level=logging.INFO)


class FileConverter:
    def __init__(self):
        self.supported_formats = {
            'pdf':  ['docx', 'txt', 'html', 'md', 'jpg', 'jpeg', 'png', 'pptx', 'ppt', 'xlsx', 'xls'],
            'docx': ['pdf'], 'doc':  ['pdf'],
            'txt':  ['pdf'], 'html': ['pdf'],
            'md':   ['pdf'], 'jpg':  ['pdf'],
            'jpeg': ['pdf'], 'png':  ['pdf'],
            'ppt':  ['pdf'], 'pptx': ['pdf'],
            'xls':  ['pdf'], 'xlsx': ['pdf'],
        }

    # ─────────────────────────────────────────────
    # HELPERS
    # ─────────────────────────────────────────────
    def _check_file(self, path):
        if not os.path.exists(path):
            raise FileNotFoundError(f"File not found: {path}")

    def _libreoffice_convert(self, input_path, output_path) -> bool:
        """Convert any Office/image file to PDF using LibreOffice headless."""
        soffice = shutil.which("soffice") or shutil.which("libreoffice")
        if not soffice:
            return False
        out_dir = os.path.dirname(os.path.abspath(output_path))
        try:
            subprocess.run(
                [soffice, "--headless", "--convert-to", "pdf",
                 "--outdir", out_dir, input_path],
                capture_output=True, text=True, timeout=120, check=True
            )
            base = os.path.splitext(os.path.basename(input_path))[0] + ".pdf"
            lo_out = os.path.join(out_dir, base)
            if os.path.exists(lo_out):
                if os.path.abspath(lo_out) != os.path.abspath(output_path):
                    shutil.move(lo_out, output_path)
                return os.path.exists(output_path) and os.path.getsize(output_path) > 0
        except Exception as e:
            logging.warning(f"LibreOffice failed: {e}")
        return False

    def _docx2pdf_convert(self, input_path, output_path) -> bool:
        """
        Convert Office file to PDF using win32com Word/WPS COM automation.
        Pixel-perfect: exact fonts, sizes, layout, images, tables.
        """
        input_abs = os.path.abspath(input_path)
        output_abs = os.path.abspath(output_path)

        # ── Strategy 1: win32com direct (most reliable on Windows) ──
        try:
            import win32com.client  # type: ignore
            import pythoncom  # type: ignore
            pythoncom.CoInitialize()
            try:
                word = win32com.client.Dispatch("Word.Application")
                word.Visible = False
                word.DisplayAlerts = False
                try:
                    doc = word.Documents.Open(input_abs, ReadOnly=True)
                    doc.SaveAs2(output_abs, FileFormat=17)  # 17 = wdFormatPDF
                    doc.Close(SaveChanges=False)
                    if os.path.exists(output_abs) and os.path.getsize(output_abs) > 0:
                        logging.info("Converted via win32com Word COM")
                        return True
                except Exception as doc_err:
                    logging.warning(f"win32com doc error: {doc_err}")
                finally:
                    try:
                        word.Quit()
                    except Exception:
                        pass
            finally:
                pythoncom.CoUninitialize()
        except Exception as e:
            logging.warning(f"win32com failed: {e}")

        # ── Strategy 2: docx2pdf wrapper ────────────────────────────
        try:
            from docx2pdf import convert  # type: ignore
            convert(input_abs, output_abs)
            if os.path.exists(output_abs) and os.path.getsize(output_abs) > 0:
                logging.info("Converted via docx2pdf")
                return True
        except Exception as e:
            logging.warning(f"docx2pdf failed: {e}")

        return False

    def _ppt_to_pdf_com(self, input_path, output_path) -> bool:
        """Convert PPT/PPTX to PDF using PowerPoint COM automation."""
        input_abs = os.path.abspath(input_path)
        output_abs = os.path.abspath(output_path)
        try:
            import win32com.client  # type: ignore
            import pythoncom  # type: ignore
            pythoncom.CoInitialize()
            try:
                ppt = win32com.client.Dispatch("PowerPoint.Application")
                ppt.Visible = True  # PowerPoint requires visible window
                try:
                    prs = ppt.Presentations.Open(input_abs, ReadOnly=True, WithWindow=False)
                    prs.SaveAs(output_abs, 32)  # 32 = ppSaveAsPDF
                    prs.Close()
                    if os.path.exists(output_abs) and os.path.getsize(output_abs) > 0:
                        logging.info("Converted PPT via win32com PowerPoint COM")
                        return True
                except Exception as e:
                    logging.warning(f"PowerPoint COM error: {e}")
                finally:
                    try:
                        ppt.Quit()
                    except Exception:
                        pass
            finally:
                pythoncom.CoUninitialize()
        except Exception as e:
            logging.warning(f"PowerPoint COM failed: {e}")
        return False

    def _excel_to_pdf_com(self, input_path, output_path) -> bool:
        """Convert XLS/XLSX to PDF using Excel COM automation."""
        input_abs = os.path.abspath(input_path)
        output_abs = os.path.abspath(output_path)
        try:
            import win32com.client  # type: ignore
            import pythoncom  # type: ignore
            pythoncom.CoInitialize()
            try:
                xl = win32com.client.Dispatch("Excel.Application")
                xl.Visible = False
                xl.DisplayAlerts = False
                try:
                    wb = xl.Workbooks.Open(input_abs, ReadOnly=True)
                    wb.ExportAsFixedFormat(0, output_abs)  # 0 = xlTypePDF
                    wb.Close(SaveChanges=False)
                    if os.path.exists(output_abs) and os.path.getsize(output_abs) > 0:
                        logging.info("Converted XLS via win32com Excel COM")
                        return True
                except Exception as e:
                    logging.warning(f"Excel COM error: {e}")
                finally:
                    try:
                        xl.Quit()
                    except Exception:
                        pass
            finally:
                pythoncom.CoUninitialize()
        except Exception as e:
            logging.warning(f"Excel COM failed: {e}")
        return False

    def _create_pdf(self, text: str, output_path: str):
        """Plain-text → PDF via fpdf2 with Unicode font."""
        try:
            from fpdf import FPDF
            pdf = FPDF(unit="pt", format="letter")
            pdf.set_auto_page_break(auto=True, margin=40)
            pdf.add_page()
            font_ok = False
            for fp in [r"C:\Windows\Fonts\arial.ttf", r"C:\Windows\Fonts\tahoma.ttf"]:
                if os.path.exists(fp):
                    try:
                        pdf.add_font("UF", style="", fname=fp)
                        pdf.set_font("UF", size=11)
                        font_ok = True
                        break
                    except Exception:
                        pass
            if not font_ok:
                pdf.set_font("Helvetica", size=11)
            for line in text.splitlines():
                if line.strip():
                    try:
                        pdf.multi_cell(0, 16, text=line)
                    except Exception:
                        pdf.multi_cell(0, 16, text=line.encode("latin-1", "replace").decode("latin-1"))
                else:
                    pdf.ln(8)
            pdf.output(output_path)
        except Exception as e:
            logging.error(f"_create_pdf fpdf2 failed: {e}")
            # reportlab fallback
            doc = SimpleDocTemplate(output_path, pagesize=letter)
            styles = getSampleStyleSheet()
            story = [Paragraph(p.strip().encode("ascii", "replace").decode("ascii"), styles["Normal"])
                     for p in text.split("\n\n") if p.strip()]
            doc.build(story)

    # ─────────────────────────────────────────────
    # PDF → OTHER
    # ─────────────────────────────────────────────
    # Font mapping: PDF font name → standard Word font
    _FONT_MAP = {
        "aptos":      "Calibri",
        "times":      "Times New Roman",
        "arial":      "Arial",
        "helvetica":  "Arial",
        "calibri":    "Calibri",
        "courier":    "Courier New",
        "verdana":    "Verdana",
        "georgia":    "Georgia",
        "cambria":    "Cambria",
        "garamond":   "Garamond",
        "tahoma":     "Tahoma",
        "trebuchet":  "Trebuchet MS",
        "palatino":   "Palatino Linotype",
        "bookman":    "Bookman Old Style",
        "century":    "Century",
        "comic":      "Comic Sans MS",
        "impact":     "Impact",
        "symbol":     "Symbol",
    }

    def _map_font(self, pdf_font_name):
        """Map a PDF font name to a standard Word font."""
        lower = pdf_font_name.lower()
        if "+" in lower:                    # strip subset prefix e.g. "ABCDEF+Aptos"
            lower = lower.split("+", 1)[1]
        for key, val in self._FONT_MAP.items():
            if key in lower:
                return val
        return "Calibri"

    def _detect_alignment(self, x0, x1, page_width, margin=50):
        """Detect paragraph alignment from line bounding box."""
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        center = (x0 + x1) / 2
        if abs(center - page_width / 2) < 40:
            return WD_ALIGN_PARAGRAPH.CENTER
        if x1 > (page_width - margin) and x0 > margin * 2:
            return WD_ALIGN_PARAGRAPH.RIGHT
        return WD_ALIGN_PARAGRAPH.LEFT

    def pdf_to_docx(self, pdf_path, output_path):
        """
        Convert PDF → DOCX with PIXEL-PERFECT layout.
        Each PDF page is rendered as a high-res image and placed full-page
        in the DOCX — guarantees identical layout, fonts, spacing, and position.
        """
        self._check_file(pdf_path)
        import fitz
        import io
        from docx.shared import Pt, Emu
        from docx.oxml.ns import qn
        from PIL import Image as PILImage

        d = docx.Document()

        # Remove default empty paragraph from blank document
        for para in d.paragraphs:
            p = para._element
            p.getparent().remove(p)

        with fitz.open(pdf_path) as pdf:
            total = len(pdf)
            for page_num, page in enumerate(pdf):
                pw_pt = page.rect.width
                ph_pt = page.rect.height
                pw_emu = int(pw_pt / 72 * 914400)
                ph_emu = int(ph_pt / 72 * 914400)

                # ── Set page size with zero margins ──────────────────────
                if page_num == 0:
                    section = d.sections[0]
                else:
                    section = d.add_section()

                section.page_width = Emu(pw_emu)
                section.page_height = Emu(ph_emu)
                section.top_margin = Emu(0)
                section.bottom_margin = Emu(0)
                section.left_margin = Emu(0)
                section.right_margin = Emu(0)
                section.header_distance = Emu(0)
                section.footer_distance = Emu(0)

                # ── Render page at 200 DPI for crisp output ───────────────
                dpi = 200
                scale = dpi / 72.0
                mat = fitz.Matrix(scale, scale)
                pix = page.get_pixmap(matrix=mat, alpha=False)

                img_bytes = io.BytesIO()
                img = PILImage.open(io.BytesIO(pix.tobytes("png")))
                img.save(img_bytes, format="PNG", optimize=False)
                img_bytes.seek(0)

                # ── Add paragraph and insert full-page image ──────────────
                para = d.add_paragraph()
                para.paragraph_format.space_before = Pt(0)
                para.paragraph_format.space_after = Pt(0)
                para.paragraph_format.line_spacing = Pt(0)

                # Set paragraph margins to zero via XML
                pPr = para._p.get_or_add_pPr()
                # Remove any indentation
                for child in list(pPr):
                    if child.tag in (qn('w:ind'), qn('w:spacing')):
                        pPr.remove(child)

                run = para.add_run()
                run.add_picture(img_bytes, width=Emu(pw_emu), height=Emu(ph_emu))

        d.save(output_path)
        logging.info(f"✅ pdf_to_docx: {total} pages rendered at {dpi}DPI → pixel-perfect DOCX")
        return True

    def pdf_to_txt(self, pdf_path, output_path):
        self._check_file(pdf_path)
        from modules.pdf_reader import extract_text_from_pdf
        text = extract_text_from_pdf(pdf_path)
        with open(output_path, "w", encoding="utf-8") as f:
            f.write(text)
        return True

    def pdf_to_html(self, pdf_path, output_path):
        """
        Convert PDF → HTML with pixel-perfect layout.
        Each page is rendered as a base64-embedded image so the HTML
        looks identical to the original PDF.
        """
        self._check_file(pdf_path)
        import fitz
        import io
        import base64
        from PIL import Image as PILImage

        pages_html = []
        with fitz.open(pdf_path) as pdf:
            for page_num, page in enumerate(pdf):
                pw_pt = page.rect.width
                ph_pt = page.rect.height
                # Render at 150 DPI — good quality, reasonable file size
                scale = 150 / 72.0
                pix = page.get_pixmap(matrix=fitz.Matrix(scale, scale), alpha=False)
                img = PILImage.open(io.BytesIO(pix.tobytes("png")))
                buf = io.BytesIO()
                img.save(buf, format="PNG", optimize=True)
                b64 = base64.b64encode(buf.getvalue()).decode("utf-8")

                # Aspect ratio for responsive display
                aspect = ph_pt / pw_pt * 100

                pages_html.append(
                    f'<div class="page" id="page-{page_num+1}">'
                    f'<div class="page-inner" style="padding-bottom:{aspect:.4f}%">'
                    f'<img src="data:image/png;base64,{b64}" alt="Page {page_num+1}" loading="lazy"/>'
                    f'</div>'
                    f'<div class="page-num">Page {page_num+1} of {len(pdf)}</div>'
                    f'</div>'
                )

        html = f"""<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="UTF-8">
<meta name="viewport" content="width=device-width, initial-scale=1.0">
<title>PDF Document</title>
<style>
  * {{ box-sizing: border-box; margin: 0; padding: 0; }}
  body {{
    background: #525659;
    font-family: Arial, sans-serif;
    padding: 20px 0;
  }}
  .page {{
    max-width: 900px;
    margin: 0 auto 24px auto;
    box-shadow: 0 4px 16px rgba(0,0,0,0.4);
    background: #fff;
  }}
  .page-inner {{
    position: relative;
    width: 100%;
    height: 0;
    overflow: hidden;
  }}
  .page-inner img {{
    position: absolute;
    top: 0; left: 0;
    width: 100%;
    height: 100%;
    display: block;
  }}
  .page-num {{
    text-align: center;
    font-size: 12px;
    color: #ccc;
    background: #525659;
    padding: 4px 0;
  }}
</style>
</head>
<body>
{''.join(pages_html)}
</body>
</html>"""

        with open(output_path, "w", encoding="utf-8") as f:
            f.write(html)
        return True

    def pdf_to_md(self, pdf_path, output_path):
        """Alias for pdf_to_markdown — fixes 'method not implemented: pdf_to_md'"""
        return self.pdf_to_markdown(pdf_path, output_path)

    def pdf_to_markdown(self, pdf_path, output_path):
        self._check_file(pdf_path)
        import fitz
        md_pages = []
        with fitz.open(pdf_path) as pdf:
            for page_num, page in enumerate(pdf):
                md_lines = [f"<!-- Page {page_num + 1} -->"]
                blocks = sorted(page.get_text("dict")["blocks"],
                                key=lambda b: (b["bbox"][1], b["bbox"][0]))
                for block in blocks:
                    if block.get("type") != 0:
                        continue
                    for line in block.get("lines", []):
                        spans = line.get("spans", [])
                        if not spans:
                            continue
                        text = " ".join(s["text"] for s in spans).strip()
                        if not text:
                            continue
                        size = spans[0].get("size", 11)
                        flags = spans[0].get("flags", 0)
                        is_bold = bool(flags & 2**4)
                        is_italic = bool(flags & 2**1)
                        if size >= 18 or (size >= 14 and is_bold):
                            md_lines.append(f"# {text}")
                        elif size >= 13 or (size >= 11 and is_bold):
                            md_lines.append(f"## {text}")
                        elif is_bold and is_italic:
                            md_lines.append(f"***{text}***")
                        elif is_bold:
                            md_lines.append(f"**{text}**")
                        elif is_italic:
                            md_lines.append(f"*{text}*")
                        else:
                            md_lines.append(text)
                    md_lines.append("")
                md_pages.append("\n".join(md_lines))
        with open(output_path, "w", encoding="utf-8") as f:
            f.write("\n\n---\n\n".join(md_pages))
        return True

    def pdf_to_jpg(self, pdf_path, output_path):
        return self._pdf_to_images(pdf_path, output_path, fmt="jpg")

    def pdf_to_jpeg(self, pdf_path, output_path):
        return self._pdf_to_images(pdf_path, output_path, fmt="jpg")

    def pdf_to_png(self, pdf_path, output_path):
        return self._pdf_to_images(pdf_path, output_path, fmt="png")

    def _pdf_to_images(self, pdf_path, output_path, fmt="jpg") -> bool:
        """
        Convert each PDF page to an image.
        If single page → saves directly to output_path.
        If multi-page  → saves as a ZIP of images.
        """
        self._check_file(pdf_path)
        import fitz
        from PIL import Image as PILImage
        import io as _io
        import zipfile

        with fitz.open(pdf_path) as doc:
            total = len(doc)
            if total == 1:
                # Single page — save directly
                pix = doc[0].get_pixmap(matrix=fitz.Matrix(2, 2))
                img = PILImage.open(_io.BytesIO(pix.tobytes("png"))).convert("RGB")
                img.save(output_path, fmt.upper() if fmt != "jpg" else "JPEG", quality=95)
            else:
                # Multi-page — zip all images
                zip_path = output_path if output_path.endswith(".zip") else output_path + ".zip"
                with zipfile.ZipFile(zip_path, "w", zipfile.ZIP_DEFLATED) as zf:
                    for i, page in enumerate(doc):
                        pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
                        img = PILImage.open(_io.BytesIO(pix.tobytes("png"))).convert("RGB")
                        buf = _io.BytesIO()
                        img.save(buf, fmt.upper() if fmt != "jpg" else "JPEG", quality=95)
                        zf.writestr(f"page_{i+1:03d}.{fmt}", buf.getvalue())
                # Point output_path to the zip
                if not output_path.endswith(".zip"):
                    import shutil as _sh
                    _sh.move(zip_path, output_path)
        return True

    def pdf_to_pptx(self, pdf_path, output_path):
        return self._pdf_to_ppt(pdf_path, output_path)

    def pdf_to_ppt(self, pdf_path, output_path):
        return self._pdf_to_ppt(pdf_path, output_path)

    def _pdf_to_ppt(self, pdf_path, output_path) -> bool:
        """
        Convert PDF to PPTX — each page becomes a slide with the page image.
        Preserves visual layout perfectly.
        """
        self._check_file(pdf_path)
        try:
            from pptx import Presentation  # type: ignore
            from pptx.util import Inches  # type: ignore
            import fitz
            import io as _io

            prs = Presentation()
            # Use blank slide layout
            blank_layout = prs.slide_layouts[6]

            with fitz.open(pdf_path) as doc:
                for page in doc:
                    # Render page at 150 DPI
                    pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
                    img_bytes = pix.tobytes("png")

                    # Set slide size to match page dimensions (in inches at 96 DPI)
                    w_in = page.rect.width / 72.0
                    h_in = page.rect.height / 72.0
                    prs.slide_width = Inches(w_in)
                    prs.slide_height = Inches(h_in)

                    slide = prs.slides.add_slide(blank_layout)
                    buf = _io.BytesIO(img_bytes)
                    slide.shapes.add_picture(buf, 0, 0,
                                             width=Inches(w_in),
                                             height=Inches(h_in))

            prs.save(output_path)
            return True
        except ImportError:
            raise Exception("python-pptx not installed. Run: pip install python-pptx")
        except Exception as e:
            logging.error(f"pdf_to_pptx error: {e}")
            raise

    def pdf_to_xlsx(self, pdf_path, output_path):
        return self._pdf_to_xls(pdf_path, output_path)

    def pdf_to_xls(self, pdf_path, output_path):
        return self._pdf_to_xls(pdf_path, output_path)

    def _pdf_to_xls(self, pdf_path, output_path) -> bool:
        """
        Convert PDF text content to Excel — each page becomes a sheet.
        Tables are detected and placed in cells.
        """
        self._check_file(pdf_path)
        try:
            import openpyxl  # type: ignore
            import fitz

            wb = openpyxl.Workbook()
            wb.remove(wb.active)  # remove default sheet

            with fitz.open(pdf_path) as doc:
                for pg_num, page in enumerate(doc):
                    ws = wb.create_sheet(title=f"Page {pg_num + 1}")
                    row_idx = 1
                    # Extract text blocks sorted by position
                    blocks = sorted(page.get_text("blocks"),
                                    key=lambda b: (round(b[1] / 20), b[0]))
                    for block in blocks:
                        text = block[4].strip()
                        if not text:
                            continue
                        # Split block into lines → each line = one row
                        for line in text.splitlines():
                            line = line.strip()
                            if not line:
                                continue
                            # Try to split by tabs or multiple spaces (table detection)
                            import re as _re
                            cells = _re.split(r'\t|  {2,}', line)
                            for col_idx, cell in enumerate(cells, start=1):
                                ws.cell(row=row_idx, column=col_idx, value=cell.strip())
                            row_idx += 1
                        row_idx += 1  # blank row between blocks

            wb.save(output_path)
            return True
        except ImportError:
            raise Exception("openpyxl not installed. Run: pip install openpyxl")
        except Exception as e:
            logging.error(f"pdf_to_xls error: {e}")
            raise

    # ─────────────────────────────────────────────
    # OTHER → PDF  (layout-preserving cascade)
    # ─────────────────────────────────────────────
    def _office_to_pdf(self, input_path, output_path) -> bool:
        """Try docx2pdf then LibreOffice for any Office file."""
        if self._docx2pdf_convert(input_path, output_path):
            logging.info(f"Converted via docx2pdf: {input_path}")
            return True
        if self._libreoffice_convert(input_path, output_path):
            logging.info(f"Converted via LibreOffice: {input_path}")
            return True
        return False

    def docx_to_pdf(self, docx_path, output_path):
        self._check_file(docx_path)
        if self._office_to_pdf(docx_path, output_path):
            return True
        # Fallback: python-docx + reportlab with bold/italic/alignment
        logging.warning("docx_to_pdf: falling back to reportlab (layout may differ)")
        doc = docx.Document(docx_path)
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
        from reportlab.lib.enums import TA_LEFT, TA_CENTER, TA_RIGHT, TA_JUSTIFY
        from reportlab.lib.pagesizes import letter as _letter
        _doc = SimpleDocTemplate(output_path, pagesize=_letter,
                                 leftMargin=72, rightMargin=72,
                                 topMargin=72, bottomMargin=72)
        styles = getSampleStyleSheet()
        align_map = {"CENTER": TA_CENTER, "RIGHT": TA_RIGHT,
                     "JUSTIFY": TA_JUSTIFY, "LEFT": TA_LEFT}
        story = []
        for para in doc.paragraphs:
            if not para.text.strip():
                story.append(Spacer(1, 6))
                continue
            sname = para.style.name if para.style else ""
            if "Heading 1" in sname:
                base = styles["Heading1"]
            elif "Heading 2" in sname:
                base = styles["Heading2"]
            elif "Heading 3" in sname:
                base = styles["Heading3"]
            else:
                base = styles["Normal"]
            text = ""
            for run in para.runs:
                t = run.text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
                if run.bold and run.italic:
                    t = f"<b><i>{t}</i></b>"
                elif run.bold:
                    t = f"<b>{t}</b>"
                elif run.italic:
                    t = f"<i>{t}</i>"
                text += t
            if not text:
                text = para.text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
            align = TA_LEFT
            if para.alignment:
                align = align_map.get(str(para.alignment).split(".")[-1], TA_LEFT)
            story.append(Paragraph(text, ParagraphStyle("c", parent=base,
                         alignment=align, spaceAfter=4)))
        _doc.build(story)
        return True

    def doc_to_pdf(self, doc_path, output_path):
        """Legacy .doc format — try Office/LibreOffice."""
        self._check_file(doc_path)
        if self._office_to_pdf(doc_path, output_path):
            return True
        raise Exception(".doc conversion requires Microsoft Word or LibreOffice")

    def txt_to_pdf(self, txt_path, output_path):
        self._check_file(txt_path)
        with open(txt_path, "r", encoding="utf-8", errors="replace") as f:
            text = f.read()
        self._create_pdf(text, output_path)
        return True

    def html_to_pdf(self, html_path, output_path):
        self._check_file(html_path)
        with open(html_path, "r", encoding="utf-8", errors="replace") as f:
            html_content = f.read()
        h = html2text.HTML2Text()
        h.ignore_links = False
        self._create_pdf(h.handle(html_content), output_path)
        return True

    def markdown_to_pdf(self, md_path, output_path):
        self._check_file(md_path)
        with open(md_path, "r", encoding="utf-8", errors="replace") as f:
            md_content = f.read()
        text = re.sub(r"<[^>]+>", "", markdown.markdown(md_content))
        self._create_pdf(text, output_path)
        return True

    def jpg_to_pdf(self, img_path, output_path):
        return self._image_to_pdf(img_path, output_path)

    def jpeg_to_pdf(self, img_path, output_path):
        return self._image_to_pdf(img_path, output_path)

    def png_to_pdf(self, img_path, output_path):
        return self._image_to_pdf(img_path, output_path)

    def _image_to_pdf(self, img_path, output_path) -> bool:
        """Convert image to PDF preserving aspect ratio."""
        self._check_file(img_path)
        try:
            import fitz
            doc = fitz.open()
            img_doc = fitz.open(img_path)
            pdfbytes = img_doc.convert_to_pdf()
            img_doc.close()
            pdf = fitz.open("pdf", pdfbytes)
            doc.insert_pdf(pdf)
            doc.save(output_path)
            doc.close()
            return True
        except Exception as e:
            logging.warning(f"fitz image→pdf failed: {e}")
        # Pillow fallback
        try:
            from PIL import Image as PILImage
            img = PILImage.open(img_path).convert("RGB")
            img.save(output_path, "PDF", resolution=150)
            return True
        except Exception as e:
            logging.error(f"Pillow image→pdf failed: {e}")
            raise

    def ppt_to_pdf(self, ppt_path, output_path):
        return self._ppt_convert(ppt_path, output_path)

    def pptx_to_pdf(self, pptx_path, output_path):
        return self._ppt_convert(pptx_path, output_path)

    def _ppt_convert(self, ppt_path, output_path) -> bool:
        self._check_file(ppt_path)
        # Try PowerPoint COM first (pixel-perfect)
        if self._ppt_to_pdf_com(ppt_path, output_path):
            return True
        # Fallback: Word COM (sometimes handles pptx)
        if self._office_to_pdf(ppt_path, output_path):
            return True
        raise Exception("PPT/PPTX conversion requires Microsoft Office or LibreOffice")

    def xls_to_pdf(self, xls_path, output_path):
        return self._xls_convert(xls_path, output_path)

    def xlsx_to_pdf(self, xlsx_path, output_path):
        return self._xls_convert(xlsx_path, output_path)

    def _xls_convert(self, xls_path, output_path) -> bool:
        self._check_file(xls_path)
        # Try Excel COM first (pixel-perfect)
        if self._excel_to_pdf_com(xls_path, output_path):
            return True
        # Fallback: LibreOffice
        if self._libreoffice_convert(xls_path, output_path):
            return True
        # Last resort: openpyxl text extraction
        try:
            import openpyxl  # type: ignore
            wb = openpyxl.load_workbook(xls_path, data_only=True)
            lines = []
            for ws in wb.worksheets:
                lines.append(f"=== {ws.title} ===")
                for row in ws.iter_rows(values_only=True):
                    lines.append("  ".join(str(c) if c is not None else "" for c in row))
            self._create_pdf("\n".join(lines), output_path)
            return True
        except Exception as e:
            logging.warning(f"openpyxl fallback failed: {e}")
        raise Exception("XLS/XLSX conversion requires Microsoft Office or LibreOffice")

    # ─────────────────────────────────────────────
    # MAIN ENTRY
    # ─────────────────────────────────────────────
    def convert_file(self, input_path, output_path, from_format, to_format):
        from_format = from_format.lower().lstrip(".")
        to_format = to_format.lower().lstrip(".")

        if from_format not in self.supported_formats:
            raise ValueError(f"Unsupported input format: {from_format}")
        if to_format not in self.supported_formats[from_format]:
            raise ValueError(f"Conversion {from_format} → {to_format} not supported")

        method = f"{from_format}_to_{to_format}"
        if hasattr(self, method):
            logging.info(f"Converting {from_format} → {to_format}")
            return getattr(self, method)(input_path, output_path)
        raise Exception(f"Conversion method not implemented: {method}")

    def get_mime_type(self, fmt):
        return {
            "pdf":  "application/pdf",
            "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            "txt":  "text/plain",
            "html": "text/html",
            "md":   "text/markdown",
        }.get(fmt, "application/octet-stream")


file_converter = FileConverter()
